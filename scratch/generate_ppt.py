import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Cấu hình đường dẫn
SCRATCH_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRATCH_DIR.parent
OUTPUT_PATH = REPO_ROOT / "Document" / "Enc2Health_Presentation.pptx"
IMAGE_PATH = REPO_ROOT / "Document" / "Enc2Health Architecture.drawio.png"

# Khởi tạo Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # tỉ lệ 16:9
prs.slide_height = Inches(7.5)

# Bảng màu chuyên nghiệp
COLOR_BG_DARK = RGBColor(15, 23, 42)      # Slate-900 (Độ sâu tối)
COLOR_TEXT_LIGHT = RGBColor(248, 250, 252) # Slate-50 (Chữ sáng)
COLOR_MUTED_LIGHT = RGBColor(148, 163, 184) # Slate-400 (Chữ mờ)

COLOR_BG_LIGHT = RGBColor(255, 255, 255)   # Trắng tinh khiết
COLOR_TEXT_DARK = RGBColor(15, 23, 42)     # Slate-900
COLOR_MUTED_DARK = RGBColor(71, 85, 105)   # Slate-600

COLOR_PRIMARY = RGBColor(139, 92, 246)     # Tím hoàng gia (Royal Purple)
COLOR_SECONDARY = RGBColor(6, 182, 212)    # Xanh ngọc (Cyan)
COLOR_ACCENT = RGBColor(16, 185, 129)      # Xanh lá (Emerald)
COLOR_DANGER = RGBColor(225, 29, 72)       # Đỏ hồng (Rose)
COLOR_BOX_BG = RGBColor(241, 245, 249)     # Slate-100 (Nền hộp nội dung)

FONT_HEADING = "Segoe UI"
FONT_BODY = "Arial"

def set_slide_background(slide, color):
    """Đặt màu nền cho slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, slide_num_text):
    """Thêm tiêu đề chuẩn cho slide nội dung."""
    # Tiêu đề slide
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_HEADING
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Số thứ tự slide ở góc phải
    txBox_num = slide.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.2), Inches(0.5))
    tf_num = txBox_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = slide_num_text
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.name = FONT_HEADING
    p_num.font.size = Pt(11)
    p_num.font.color.rgb = COLOR_MUTED_DARK

    # Đường kẻ phân cách mỏng
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(12.13), Inches(0.02)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = RGBColor(226, 232, 240) # Slate-200
    connector.line.fill.background()

def create_info_box(slide, left, top, width, height, title, body, border_color=COLOR_PRIMARY, bg_color=COLOR_BOX_BG):
    """Tạo một hộp nội dung có viền màu làm điểm nhấn (giống Highlight Box trong HTML)."""
    # Tạo hộp nền
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(1.5)
    
    # Tạo text frame bên trong
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = tf.margin_bottom = Inches(0.15)
    
    # Paragraph tiêu đề hộp
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = border_color
    p_title.space_after = Pt(6)
    
    # Paragraph nội dung
    p_body = tf.add_paragraph()
    p_body.text = body
    p_body.font.name = FONT_BODY
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = COLOR_TEXT_DARK
    p_body.line_spacing = 1.2
    
    return box

# ==============================================================================
# SLIDE 1: Slide Tiêu Đề (Dark Background)
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank slide
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1, COLOR_BG_DARK)

# Text box chính giữa
tx_title = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
tf1 = tx_title.text_frame
tf1.word_wrap = True

# Logo/Icon
p_icon = tf1.paragraphs[0]
p_icon.text = "🔐"
p_icon.font.size = Pt(54)
p_icon.alignment = PP_ALIGN.CENTER
p_icon.space_after = Pt(10)

# Tiêu đề
p_main = tf1.add_paragraph()
p_main.text = "Enc²Health"
p_main.font.name = FONT_HEADING
p_main.font.size = Pt(56)
p_main.font.bold = True
p_main.font.color.rgb = COLOR_TEXT_LIGHT
p_main.alignment = PP_ALIGN.CENTER
p_main.space_after = Pt(8)

# Phụ đề
p_sub = tf1.add_paragraph()
p_sub.text = "Hybrid Adaptive Encrypted Query Processing cho Cloud-Native DBMS"
p_sub.font.name = FONT_HEADING
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_SECONDARY
p_sub.alignment = PP_ALIGN.CENTER
p_sub.space_after = Pt(2)

p_sub2 = tf1.add_paragraph()
p_sub2.text = "Bảo vệ dữ liệu y tế nhạy cảm tự thích ứng trong môi trường đám mây"
p_sub2.font.name = FONT_HEADING
p_sub2.font.size = Pt(14)
p_sub2.font.color.rgb = COLOR_MUTED_LIGHT
p_sub2.alignment = PP_ALIGN.CENTER

# Thông tin tác giả ở dưới cùng
tx_authors = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(1.0))
tf_auth = tx_authors.text_frame
tf_auth.word_wrap = True
p_auth = tf_auth.paragraphs[0]
p_auth.text = "Nguyễn Hoàng Long (KMS & Enc)   |   Lâm Tú Lan (TEE & Obs)   |   Nguyễn Lê Thành Nam (Router & Adaptive)"
p_auth.font.name = FONT_HEADING
p_auth.font.size = Pt(12)
p_auth.font.color.rgb = COLOR_MUTED_LIGHT
p_auth.alignment = PP_ALIGN.CENTER

# ==============================================================================
# SLIDE 2: Thực Trạng & Rủi Ro Bảo Mật (Light Background)
# ==============================================================================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2, COLOR_BG_LIGHT)
add_header(slide2, "🚨 Thực Trạng & Rủi Ro Bảo Mật (Security Risks)", "Slide 2 / 14")

# Cột Trái: Rủi ro bảo mật chính
create_info_box(
    slide2, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.2),
    "Thách thức: Lộ lọt dữ liệu y tế nhạy cảm (PHI)",
    "Thông tin y tế nhạy cảm khi lưu trữ đám mây dễ bị xâm phạm do Cloud Admin tò mò (mô hình đe dọa Honest-but-curious) hoặc Hacker chiếm quyền ảo hóa OS/Hypervisor.\n\n"
    "• Rò rỉ tĩnh (At-Rest): Đọc trộm file DB vật lý trên ổ cứng đám mây.\n"
    "• Rò rỉ tính toán (In-Use): Lớn nhất. DBMS giải mã ciphertext thành plaintext trên RAM của host để tính toán gộp (SUM/AVG) làm lộ bộ nhớ.",
    border_color=COLOR_DANGER
)

# Cột Phải: Sự kiện thực tế chứng minh
create_info_box(
    slide2, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "Sự Kiện Thực Tế & Rủi Ro Thực Tiễn",
    "• Vụ tấn công SingHealth (Singapore - 2018):\n"
    "  Đánh cắp bệnh án của 1.5 triệu bệnh nhân, bao gồm cả Thủ tướng Lý Hiển Long, do kiểm soát DB đám mây.\n\n"
    "• Vụ Change Healthcare (Mỹ - 2024):\n"
    "  Rò rỉ thông tin y tế của hơn 1/3 người dân Mỹ trên cloud, gây tê liệt thanh toán và dịch vụ y tế diện rộng.\n\n"
    "• Sự cố rò rỉ RAM (Meltdown/Spectre & VM Escape):\n"
    "  Kẻ tấn công hoặc Cloud admin bất chính vượt qua ranh giới ảo hóa của máy chủ đám mây để đọc trực tiếp dữ liệu và khóa giải mã plaintext trên RAM của máy chủ.",
    border_color=COLOR_MUTED_DARK
)

# Thêm hộp mô hình đe dọa góc dưới bên trái
create_info_box(
    slide2, Inches(0.6), Inches(3.9), Inches(5.8), Inches(2.9),
    "Mô hình đe dọa giả định (Threat Model)",
    "Hệ thống giả định nhà cung cấp đám mây là Honest-but-curious:\n"
    "1. Vận hành hệ thống đúng cách (trung thực).\n"
    "2. Nhưng tò mò: Cố gắng thu thập thông tin y tế bệnh nhân bằng cách quét đĩa cứng hoặc dump bộ nhớ RAM khi DBMS tính toán.",
    border_color=COLOR_SECONDARY
)

# ==============================================================================
# SLIDE 3: Mục Tiêu Của Đồ Án (Light Background)
# ==============================================================================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3, COLOR_BG_LIGHT)
add_header(slide3, "🎯 Mục Tiêu Của Đồ Án (Project Goals)", "Slide 3 / 14")

# 3 hộp mục tiêu song song
create_info_box(
    slide3, Inches(0.6), Inches(1.8), Inches(3.8), Inches(4.0),
    "1. An Toàn Khi Tính Toán",
    "Đảm bảo dữ liệu y tế nhạy cảm (PHI) luôn được bảo vệ dưới dạng bản mã khi lưu trữ trên cơ sở dữ liệu và chỉ được giải mã tính toán bên trong phân vùng nhớ mã hóa phần cứng (Enclave TEE).\n\n"
    "Không để lộ plaintext hoặc khóa giải mã trên RAM thường của host đám mây.",
    border_color=COLOR_PRIMARY
)

create_info_box(
    slide3, Inches(4.7), Inches(1.8), Inches(3.8), Inches(4.0),
    "2. Tính Tương Thích & Tiện Dụng",
    "Hỗ trợ đầy đủ các toán tử phân tích và thống kê y tế phổ biến:\n\n"
    "• Lọc so khớp chính xác (=, DTE)\n"
    "• Tìm kiếm khoảng tuổi (ORE)\n"
    "• Tính toán gộp nhạy cảm (SUM, AVG)\n"
    "• Tìm kiếm triệu chứng lâm sàng bằng từ khóa (Static SSE index).",
    border_color=COLOR_SECONDARY
)

create_info_box(
    slide3, Inches(8.8), Inches(1.8), Inches(3.9), Inches(4.0),
    "3. Sẵn Sàng Tự Thích Ứng",
    "Khắc phục điểm yếu giới hạn dung lượng phần cứng bảo mật (EPC Enclave).\n\n"
    "Tự động điều tiết tải và chuyển đổi định tuyến truy vấn linh hoạt (Fallback Mode) khi Enclave chịu áp lực tài nguyên lớn, đảm bảo hệ thống luôn sẵn sàng 24/7.",
    border_color=COLOR_ACCENT
)

# Trích dẫn dưới cùng
tx_quote = slide3.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.33), Inches(0.8))
tf_q = tx_quote.text_frame
p_q = tf_q.paragraphs[0]
p_q.text = "\"Đạt được sự cân bằng tối ưu giữa bảo mật dữ liệu y tế tuyệt đối và hiệu năng/tính sẵn sàng của hệ thống\""
p_q.font.name = FONT_HEADING
p_q.font.size = Pt(14)
p_q.font.italic = True
p_q.font.color.rgb = COLOR_MUTED_DARK
p_q.alignment = PP_ALIGN.CENTER

# ==============================================================================
# SLIDE 4: Kiến Trúc Giải Pháp Lai Tự Thích Nghi (Light Background)
# ==============================================================================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4, COLOR_BG_LIGHT)
add_header(slide4, "🏗️ Kiến Trúc Giải Pháp Lai Tự Thích Nghi", "Slide 4 / 14")

# Cột Trái: Chi tiết các thành phần chế độ
create_info_box(
    slide4, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "Cơ Chế Định Tuyến Lai Hướng Toán Tử",
    "• Software Mode (Xử lý trên Bản mã):\n"
    "  Thực thi lọc và đếm cơ bản (=, range, COUNT) trực tiếp trên bản mã lưu tại MongoDB bằng mật mã học DTE và ORE.\n"
    "  Máy chủ đám mây tính toán mà không cần biết nội dung plaintext.\n\n"
    "• TEE Mode (Tính toán trong Enclave):\n"
    "  Chuyển các toán tử nhạy cảm đòi hỏi giải mã (SUM, AVG, COUNT DISTINCT) vào vùng nhớ mã hóa phần cứng của Intel SGX Enclave để xử lý thông qua DuckDB in-memory.\n\n"
    "• Adaptive Switching (Tự thích nghi):\n"
    "  Router liên tục giám sát tải của Enclave. Nếu áp lực EPC bão hòa (>=80%), Router tự động hạ cấp toán tử TEE về chạy ở chế độ Software Fallback để đảm bảo không bị nghẽn luồng dịch vụ.",
    border_color=COLOR_PRIMARY
)

# Cột Phải: Sơ đồ dòng quyết định
create_info_box(
    slide4, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "Nguyên Lý Quyết Định Định Tuyến",
    "• Giao diện Client gửi câu lệnh kèm JWT xác thực người dùng.\n\n"
    "• Router phân tích toán tử:\n"
    "  - Nếu là toán tử an toàn (COUNT, SEARCH, =):\n"
    "    Định tuyến chạy Software Mode trên MongoDB ciphertext.\n"
    "  - Nếu là toán tử thống kê nhạy cảm (SUM, AVG):\n"
    "    Đọc trạng thái EPC Enclave:\n"
    "      + Nếu EPC khả dụng (<80%): Định tuyến TEE Mode.\n"
    "      + Nếu EPC quá tải (>=80%): Kích hoạt Software Fallback (giải mã an toàn cục bộ trên host bằng DEK).\n\n"
    "• Kết quả trả về được Router lọc quyền RBAC/ABAC trước khi hiển thị lên trình duyệt của Client.",
    border_color=COLOR_SECONDARY
)

# ==============================================================================
# SLIDE 5: Sơ Đồ Kiến Trúc Hệ Thống (Light Background + Nhúng Ảnh)
# ==============================================================================
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5, COLOR_BG_LIGHT)
add_header(slide5, "🏗️ Sơ Đồ Kiến Trúc Hệ Thống (Enc²Health Architecture)", "Slide 5 / 14")

if IMAGE_PATH.exists():
    # Nhúng hình ảnh sơ đồ kiến trúc draw.io
    slide5.shapes.add_picture(
        str(IMAGE_PATH),
        Inches(1.5), Inches(1.5), Inches(10.33), Inches(4.8)
    )
    
    # Caption hình ảnh
    tx_caption = slide5.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(11.33), Inches(0.6))
    tf_c = tx_caption.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = "* Sơ đồ chi tiết thiết kế hệ thống 5 phân lớp bảo mật liên thông (kết xuất từ bản vẽ Enc2Health_Architecture.xml)"
    p_c.font.name = FONT_HEADING
    p_c.font.size = Pt(11)
    p_c.font.color.rgb = COLOR_MUTED_DARK
    p_c.alignment = PP_ALIGN.CENTER
else:
    # Nếu không tìm thấy hình ảnh, fallback sang text box giải thích cấu trúc
    create_info_box(
        slide5, Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0),
        "Sơ đồ 5 lớp bảo mật liên thông",
        "1. Client Environment: Gửi request qua giao diện Web HTTP kèm JWT.\n"
        "2. Query Router Layer (:8000): Điều khiển logic, xác thực JWT, RBAC/ABAC, Adaptive Controller.\n"
        "3. Cloud-Native DBMS (:27017): Lưu trữ patient_records và sse_index đã mã hóa trên MongoDB.\n"
        "4. TEE Intel SGX Enclave (:9091): Môi trường giải mã an toàn và tính toán phân tích DuckDB.\n"
        "5. Key Management Service (:8200): Quản lý khóa bọc tập trung sử dụng HashiCorp Vault.",
        border_color=COLOR_PRIMARY
    )

# ==============================================================================
# SLIDE 6: Tầng Mật Mã Học Chọn Lọc (Light Background)
# ==============================================================================
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6, COLOR_BG_LIGHT)
add_header(slide6, "🔑 Tầng Mật Mã Học Chọn Lọc (Cryptographic Layers)", "Slide 6 / 14")

# Chia thành 2 cột hiển thị các nhóm mật mã
create_info_box(
    slide6, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "1. Nhóm mật mã hóa chạy trên host (Software Mode)",
    "• DTE (Deterministic Encryption - AES-SIV-256):\n"
    "  - Mã hóa trường: ma_benh, khoa_phong, CCCD.\n"
    "  - Tính chất: Cùng giá trị gốc cho ra cùng bản mã. Cho phép MongoDB xây dựng chỉ mục và tìm kiếm bằng chính xác trực tiếp.\n\n"
    "• ORE/OPE (Order-Preserving - Boldyreva):\n"
    "  - Mã hóa trường: tuoi, ngay_nhap_vien.\n"
    "  - Tính chất: Bảo toàn thứ tự toán học của plaintext. Cho phép tìm kiếm khoảng (>, <, between) và lập chỉ mục cây B-Tree.\n\n"
    "• Static SSE (Searchable Symmetric Encryption):\n"
    "  - Mã hóa trường: triệu chứng lâm sàng dạng văn bản dài.\n"
    "  - Tính chất: Dùng HMAC băm từ khóa + postings mã hóa AES-GCM.",
    border_color=COLOR_PRIMARY
)

create_info_box(
    slide6, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "2. Nhóm mật mã hóa an toàn (TEE Mode / KMS)",
    "• AES-GCM-256 (Mã hóa đối xứng ngẫu nhiên):\n"
    "  - Mã hóa trường: vien_phi, ket_qua_xet_nghiem.\n"
    "  - Tính chất: Bản mã ngẫu nhiên tuyệt đối, tích hợp xác thực toàn vẹn AEAD. Dữ liệu y tế tối mật chỉ được giải mã và tính toán gộp (DuckDB) bên trong Enclave TEE.\n\n"
    "• ECC ECIES P-384 (Mã hóa bất đối xứng):\n"
    "  - Mã hóa trường: pii_enc (Họ tên, CCCD, Bệnh án chi tiết).\n"
    "  - Tính chất: Phía Client mã hóa bằng Public Key của từng khoa. Private Key chỉ được nạp từ Vault vào Enclave TEE để giải mã.\n\n"
    "• Envelope Encryption (Mã hóa phong bì):\n"
    "  - Khóa DEK (gcm_dek) được mã hóa bọc bởi Master Key trong HashiCorp Vault KV-v2 và chỉ giải bọc về RAM Enclave.",
    border_color=COLOR_SECONDARY
)

# ==============================================================================
# SLIDE 7: Cơ Chế Tự Thích Nghi & Hysteresis (Light Background)
# ==============================================================================
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7, COLOR_BG_LIGHT)
add_header(slide7, "⚡ Cơ Chế Tự Thích Nghi & Hysteresis (Adaptive Logic)", "Slide 7 / 14")

# Cột Trái: Trạng thái bão hòa EPC và fallback
create_info_box(
    slide7, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "EPC Saturation & Fallback Hysteresis",
    "• Vấn đề dao động liên tục (Flapping):\n"
    "  Nếu chỉ dùng một ngưỡng duy nhất (ví dụ 80% áp lực bộ nhớ) để chuyển đổi trạng thái, khi tải dao động nhẹ quanh mức 80%, hệ thống sẽ liên tục bật/tắt chế độ fallback. Việc này làm mất ổn định hệ thống.\n\n"
    "• Giải pháp Hysteresis (Vùng trễ hai ngưỡng):\n"
    "  - Ngưỡng Kích hoạt Fallback: EPC Pressure >= 80%.\n"
    "  - Ngưỡng Khôi phục TEE: EPC Pressure <= 60%.\n"
    "  - Vùng trễ (60% - 80%): Giữ nguyên chế độ hoạt động hiện tại để đảm bảo trạng thái hệ thống ổn định lâu dài.\n\n"
    "• Software Fallback:\n"
    "  Khi quá tải, Router tạm thời giải mã bằng khóa cục bộ trên RAM host để xử lý phép tính, ưu tiên tính sẵn sàng (Availability) của bệnh viện.",
    border_color=COLOR_DANGER
)

# Cột Phải: Prober và Resource Monitor
create_info_box(
    slide7, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "Cơ Chế Giám Sát Chủ Động (EPC Prober)",
    "• Hoạt động định kỳ:\n"
    "  Một luồng chạy nền trong Router gửi truy vấn COUNT nhỏ đến Enclave Pool mỗi 5 giây để đo độ trễ latency.\n\n"
    "• Khóa baseline tự động:\n"
    "  Trong 3 lần chạy đầu tiên ở trạng thái rảnh rỗi, hệ thống tính toán và khóa lại giá trị Baseline Latency làm đường cơ sở.\n\n"
    "• Phát hiện bão hòa:\n"
    "  Nếu latency phép thử tăng gấp đôi so với Baseline (latency ratio >= 2.0), prober xác định EPC đang bão hòa.\n\n"
    "• Kết hợp số liệu thực tế:\n"
    "  Router lấy thông số bộ nhớ RSS thực tế của tiến trình Pool qua đường dẫn hệ thống /proc/<pid>/status để làm căn cứ chính xác nhất.",
    border_color=COLOR_SECONDARY
)

# ==============================================================================
# SLIDE 8: Phân Quyền & Kiểm Soát Truy Cập ABAC / RBAC (Light Background)
# ==============================================================================
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8, COLOR_BG_LIGHT)
add_header(slide8, "🛡️ Phân Quyền & Kiểm Soát Truy Cập ABAC / RBAC", "Slide 8 / 14")

# Cột Trái: RBAC Column-level
create_info_box(
    slide8, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "RBAC: Phân quyền theo Cột (Column-Level)",
    "Hệ thống kiểm tra JWT HS256 gửi kèm từ phía Client để phân loại vai trò và che giấu trường thông tin nhạy cảm:\n\n"
    "• Bác sĩ (doctor) & Quản trị (admin):\n"
    "  Được quyền xem đầy đủ thông tin bệnh án và phác đồ điều trị sau khi Enclave giải mã.\n\n"
    "• Nhân viên hành chính (admin_staff):\n"
    "  Làm thủ tục tiếp đón, được xem thông tin cá nhân plaintext (Họ tên, CCCD, địa chỉ, ngày sinh) nhưng hai trường lâm sàng nhạy cảm bị che giấu cứng thành [MASKED] (Tránh lộ thông tin bệnh án).\n\n"
    "• Nhà nghiên cứu (researcher):\n"
    "  Chỉ được thống kê gộp, các trường viện phí và mã bệnh bị che giấu thành [MASKED] để bảo vệ thông tin.",
    border_color=COLOR_PRIMARY
)

# Cột Phải: ABAC Dept-scoping
create_info_box(
    slide8, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "ABAC: Giới Hạn Khoa Phòng (Dept-Scoping)",
    "Ngăn chặn việc truy cập ngang trái phép hồ sơ bệnh án giữa các khoa phòng y tế:\n\n"
    "• Chính sách: Bác sĩ điều trị thuộc khoa nào chỉ được truy vấn thông tin của bệnh nhân thuộc khoa đó.\n\n"
    "• Cơ chế Zero-Trust Client:\n"
    "  - Router trích xuất thuộc tính bộ phận (dept) từ các claims của chữ ký JWT.\n"
    "  - Router tự động mã hóa giá trị khoa bằng DTE khoa và tiêm filter khoa phòng bắt buộc vào truy vấn gửi lên MongoDB.\n"
    "  - Client không thể can thiệp hay xóa bỏ filter khoa phòng y tế này.\n"
    "  - Phản ứng: Nếu bác sĩ cố tình sửa ID bệnh nhân khoa khác để tra cứu, Router kiểm tra thuộc tính khoa phòng và trả về 403 Forbidden.",
    border_color=COLOR_ACCENT
)

# ==============================================================================
# SLIDE 9: Phương Án Triển Khai (Light Background)
# ==============================================================================
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9, COLOR_BG_LIGHT)
add_header(slide9, "⚙️ Phương Án Triển Khai (Deployment Strategy)", "Slide 9 / 14")

# Cột Trái: Công nghệ và đóng gói
create_info_box(
    slide9, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "Tích Hợp TEE Gramine & HashiCorp Vault",
    "• Đóng gói Gramine LibOS:\n"
    "  - Xây dựng manifest cho DuckDB và runtime Python.\n"
    "  - Đo lường dung lượng RAM tối thiểu (128MB) và ký số MRENCLAVE cho các tệp tin cấu trúc tin cậy.\n"
    "  - DuckDB chạy trực tiếp trong TEE mà không cần sửa đổi mã nguồn C/C++.\n\n"
    "• Bảo mật KMS qua Vault AppRole:\n"
    "  - Khi Enclave startup, gửi yêu cầu chứng thực lấy token tạm thời từ Vault thông qua Role_ID và Secret_ID.\n"
    "  - Đọc khóa giải mã y tế (gcm_dek) và private key ECC của khoa về vùng RAM bảo mật của Enclave.\n"
    "  - Không lưu khóa plaintext ra đĩa cứng của đám mây.\n\n"
    "• mTLS (Mutual TLS):\n"
    "  Sử dụng chứng chỉ CA nội bộ để thiết lập mã hóa đường truyền an toàn cho luồng dữ liệu đi từ Enclave Pool về Router.",
    border_color=COLOR_PRIMARY
)

# Cột Phải: Các lệnh vận hành
create_info_box(
    slide9, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "Lệnh Vận Hành & Smoke Test Tích Hợp",
    "• Khởi chạy toàn bộ hệ thống (One-command startup):\n"
    "  $ make up\n"
    "  (Khởi động MongoDB, seeder gieo 10,000 hồ sơ EHR ciphertext, khởi động Vault, ECALL Task Pool và Query Router)\n\n"
    "• Lệnh phụ trợ:\n"
    "  - Xem trạng thái: make ps\n"
    "  - Xem log: make logs\n"
    "  - Tắt hệ thống và xóa volume: make down\n\n"
    "• Khởi chạy Smoke Test mTLS E2E nội bộ:\n"
    "  $ make smoke-local\n\n"
    "• Kết quả E2E test xác thực liên thông:\n"
    "  $ pytest tests/test_e2e.py -v\n"
    "  => 10/10 Passed (Thống kê và PII chạy thực tế trên MongoDB, Vault và Enclave Pool).",
    border_color=COLOR_SECONDARY
)

# ==============================================================================
# SLIDE 10: Kết Quả Đánh Giá Hiệu Năng (Light Background)
# ==============================================================================
slide10 = prs.slides.add_slide(slide_layout)
set_slide_background(slide10, COLOR_BG_LIGHT)
add_header(slide10, "📊 Kết Quả Đánh Giá Hiệu Năng (Performance)", "Slide 10 / 14")

# Tạo 2 bảng/hộp nội dung song song để hiển thị benchmark
create_info_box(
    slide10, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "1. Thời gian phản hồi (50 runs mỗi chế độ)",
    "• SOFTWARE (count y tế):\n"
    "  - Độ trễ trung bình: 5.7 ms\n"
    "  - p95: 6.2 ms | QPS: 174\n\n"
    "• TEE (sum_vien_phi y tế):\n"
    "  - Độ trễ trung bình: 10.6 ms\n"
    "  - p95: 11.6 ms | QPS: 94\n\n"
    "• TEE (avg_vien_phi y tế):\n"
    "  - Độ trễ trung bình: 11.0 ms\n"
    "  - p95: 12.0 ms | QPS: 91\n\n"
    "• Hybrid + RBAC masking y tế:\n"
    "  - Độ trễ trung bình: 11.8 ms\n"
    "  - p95: 15.1 ms | QPS: 84\n\n"
    "=> Nhận xét: TEE overhead chỉ cao gấp ~2 lần so với Software. Đây là độ trễ tối ưu và hoàn toàn có thể chấp nhận được trong ngành y tế.",
    border_color=COLOR_PRIMARY
)

create_info_box(
    slide10, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "2. Khả năng chịu tải (Concurrent Clients)",
    "Hệ thống kiểm thử với số lượng yêu cầu truy vấn đồng thời từ 1 đến 50 clients kết nối vào API thống kê viện phí trung bình:\n\n"
    "• 1 Client đồng thời:\n"
    "  - Độ trễ trung bình: 7.1 ms | QPS: 141 | Lỗi: 0%\n\n"
    "• 10 Clients đồng thời:\n"
    "  - Độ trễ trung bình: 31.6 ms | QPS: 305 | Lỗi: 0%\n\n"
    "• 50 Clients đồng thời:\n"
    "  - Độ trễ trung bình: 152.9 ms | QPS: 299 | Lỗi: 0%\n\n"
    "=> Nhận xét: Nhờ sử dụng Thread Pool 8 workers và tập lệnh AES-NI phần cứng, hệ thống duy trì được throughput ổn định (khoảng 300 QPS) với tỷ lệ lỗi 0% tuyệt đối ở mức tải cao.",
    border_color=COLOR_SECONDARY
)

# ==============================================================================
# SLIDE 11: Đánh Giá An Ninh & Rò Rỉ (Light Background)
# ==============================================================================
slide11 = prs.slides.add_slide(slide_layout)
set_slide_background(slide11, COLOR_BG_LIGHT)
add_header(slide11, "🔒 Đánh Giá An Ninh & Rò Rỉ (Security Assessment)", "Slide 11 / 14")

# Cột Trái: Q-leakage và Software Fallback
create_info_box(
    slide11, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "Phân Tích Q-Leakage khi Fallback",
    "• Trạng thái TEE Mode (Bình thường):\n"
    "  - Rò rỉ dữ liệu = 0.\n"
    "  - Các kết quả được tổng hợp trực tiếp trong Enclave. Dữ liệu trung gian và khóa DEK được cô lập hoàn toàn trước Hypervisor.\n\n"
    "• Trạng thái Software Fallback Mode:\n"
    "  - Khi Enclave quá tải bộ nhớ, hệ thống chuyển sang giải mã cục bộ bằng CPU host để tính toán gộp.\n"
    "  - Khóa giải mã DEK (gcm_dek) và các bản ghi plaintext xuất hiện trên RAM host tạo ra rủi ro bị rò rỉ nếu Hypervisor bị dump RAM.\n\n"
    "• Đánh giá mức độ lộ lọt dữ liệu (Q-Leakage):\n"
    "  Tệp tests/leakage.py thực hiện đo đạc sự khác biệt rõ rệt về mức độ lộ lọt dữ liệu giữa TEE mode (đã masked) so với Software Fallback raw.",
    border_color=COLOR_DANGER
)

# Cột Phải: ORE & SSE Leakage
create_info_box(
    slide11, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "Đánh Giá Bảo Mật Thuật Toán ORE & SSE",
    "• ORE Order-Leakage (Tuổi bệnh nhân):\n"
    "  - Thực nghiệm đo lường sai số phục hồi tuổi thực tế (MAE - Mean Absolute Error) khi kẻ tấn công thực hiện tấn công ánh xạ thứ hạng là 0.67 tuổi (cho kịch bản có neo).\n"
    "  - Hệ thống hạn chế rủi ro bằng cách chỉ dùng ORE lọc thô tuổi ở Software mode.\n\n"
    "• Static SSE Leakage (Triệu chứng):\n"
    "  - Rò rỉ search pattern (HMAC token trùng nhau) và access pattern (postings y tế lộ ID bệnh nhân trùng).\n"
    "  - Khắc phục bằng cách tách biệt chỉ mục sse_index với các thông tin định danh cá nhân y tế nhạy cảm (PII).",
    border_color=COLOR_SECONDARY
)

# ==============================================================================
# SLIDE 12: Kịch Bản Demo Thực Tế (Light Background)
# ==============================================================================
slide12 = prs.slides.add_slide(slide_layout)
set_slide_background(slide12, COLOR_BG_LIGHT)
add_header(slide12, "🎬 Kịch Bản Demo Thực Tế (Live Demo Walkthrough)", "Slide 12 / 14")

# 3 kịch bản demo
create_info_box(
    slide12, Inches(0.6), Inches(1.8), Inches(3.8), Inches(4.0),
    "1. ABAC Dept-Scoping",
    "• Chạy script:\n"
    "  demo_abac.py\n\n"
    "• Diễn biến:\n"
    "  Bác sĩ khoa Tim mạch chỉ truy vấn được 1,700 hồ sơ bệnh nhân thuộc khoa mình. Sửa filter tìm kiếm khoa khác lập tức bị Router phát hiện và trả về lỗi 403 Forbidden.\n"
    "  (Chứng minh Zero-trust client)",
    border_color=COLOR_SECONDARY
)

create_info_box(
    slide12, Inches(4.7), Inches(1.8), Inches(3.8), Inches(4.0),
    "2. Staff PII Masking",
    "• Chạy script:\n"
    "  demo_e2e.py (luồng PII)\n\n"
    "• Diễn biến:\n"
    "  Nhân viên hành chính tra cứu hồ sơ bệnh nhân bằng CCCD: thấy Họ tên, CCCD plaintext; nhưng hai trường y tế nhạy cảm là \"Tóm tắt bệnh án\" và \"Phác đồ điều trị\" bị Router che giấu thành [MASKED].",
    border_color=COLOR_ACCENT
)

create_info_box(
    slide12, Inches(8.8), Inches(1.8), Inches(3.9), Inches(4.0),
    "3. Live Fallback & Restore",
    "• Chạy script:\n"
    "  demo_adaptive.py\n\n"
    "• Diễn biến:\n"
    "  Mô phỏng đợt dịch làm tăng tải. Ép áp lực EPC lên 90% via API. Router tự fallback các truy vấn nhạy cảm từ TEE sang Software.\n"
    "  Giải tỏa áp lực xuống 30%: Router khôi phục về chế độ TEE (Hysteresis pass).",
    border_color=COLOR_DANGER
)

# ==============================================================================
# SLIDE 13: Hạn Chế & Hướng Phát Triển (Light Background)
# ==============================================================================
slide13 = prs.slides.add_slide(slide_layout)
set_slide_background(slide13, COLOR_BG_LIGHT)
add_header(slide13, "⚠️ Hạn Chế & Hướng Phát Triển (Future Work)", "Slide 13 / 14")

# Hộp Hạn chế
create_info_box(
    slide13, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3),
    "Các Hạn Chế Hiện Tại",
    "• Giả lập Gramine Simulation:\n"
    "  Hệ thống chạy ở chế độ simulation, vùng nhớ EPC chưa thực sự được mã hóa vật lý bằng phần cứng Intel SGX.\n\n"
    "• Chứng thực Attestation đơn giản:\n"
    "  Sử dụng mô phỏng chữ ký HMAC + freshness để xác thực enclave, chưa phải là SGX Quote / DCAP thật.\n\n"
    "• Rò rỉ khóa khi Software Fallback:\n"
    "  Khi quá tải, việc nạp khóa giải mã gcm_dek ra RAM host để tính toán là một sự đánh đổi về mặt an ninh để giữ tính sẵn sàng.",
    border_color=COLOR_DANGER
)

# Hộp Hướng phát triển
create_info_box(
    slide13, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
    "Hướng Phát Triển Tiếp Theo",
    "• Triển khai phần cứng vật lý thật:\n"
    "  Đưa hệ thống lên đám mây hỗ trợ Intel SGX2 vật lý (ví dụ: Azure DCsv3) và cấu hình cơ chế chứng thực **RA-TLS (Remote Attestation TLS)**.\n\n"
    "• Áp dụng Mã Hóa Đồng Cấu (HE):\n"
    "  Nghiên cứu áp dụng thuật toán **mã hóa đồng cấu Paillier** cho cơ chế Fallback. Máy chủ host có thể thực hiện phép tính tổng trực tiếp trên bản mã mà không cần giải mã dữ liệu hay giữ khóa giải mã trên RAM host.\n\n"
    "• Che giấu access pattern:\n"
    "  Tích hợp ORAM (Oblivious RAM) và chèn dữ liệu nhiễu để nâng cao độ bảo mật cho ORE và SSE.",
    border_color=COLOR_ACCENT
)

# ==============================================================================
# SLIDE 14: Slide Kết Luận & Hỏi Đáp (Dark Background)
# ==============================================================================
slide14 = prs.slides.add_slide(slide_layout)
set_slide_background(slide14, COLOR_BG_DARK)

# Nội dung chính
tx_end = slide14.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
tf_end = tx_end.text_frame
tf_end.word_wrap = True

# Logo cuối
p_end_icon = tf_end.paragraphs[0]
p_end_icon.text = "💡"
p_end_icon.font.size = Pt(54)
p_end_icon.alignment = PP_ALIGN.CENTER
p_end_icon.space_after = Pt(10)

# Tiêu đề cảm ơn
p_cảm_ơn = tf_end.add_paragraph()
p_cảm_ơn.text = "Xin Cảm Ơn Hội Đồng Đã Lắng Nghe!"
p_cảm_ơn.font.name = FONT_HEADING
p_cảm_ơn.font.size = Pt(44)
p_cảm_ơn.font.bold = True
p_cảm_ơn.font.color.rgb = COLOR_TEXT_LIGHT
p_cảm_ơn.alignment = PP_ALIGN.CENTER
p_cảm_ơn.space_after = Pt(10)

# Lời dẫn
p_end_lead = tf_end.add_paragraph()
p_end_lead.text = "Enc²Health — Giải pháp bảo vệ dữ liệu y tế lai tự thích nghi trên Cloud-Native DBMS"
p_end_lead.font.name = FONT_HEADING
p_end_lead.font.size = Pt(18)
p_end_lead.font.color.rgb = COLOR_SECONDARY
p_end_lead.alignment = PP_ALIGN.CENTER
p_end_lead.space_after = Pt(15)

p_qa = tf_end.add_paragraph()
p_qa.text = "Q&A — RẤT MONG NHẬN ĐƯỢC CÂU HỎI VÀ ĐÓNG GÓP TỪ CÁC THẦY CÔ"
p_qa.font.name = FONT_HEADING
p_qa.font.size = Pt(14)
p_qa.font.bold = True
p_qa.font.color.rgb = COLOR_ACCENT
p_qa.alignment = PP_ALIGN.CENTER

# Lưu file
prs.save(str(OUTPUT_PATH))
print(f"Presentation saved successfully to: {OUTPUT_PATH}")
